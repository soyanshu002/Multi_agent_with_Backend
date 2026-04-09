from pathlib import Path

import pandas as pd
from langchain_community.document_loaders import CSVLoader, PyPDFLoader
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from docx import Document as DocxDocument
from pptx import Presentation


SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".csv",
    ".docx",
    ".pptx",
    ".xlsx",
    ".xls",
    ".txt",
    ".md",
    ".rtf",
    ".json",
    ".html",
    ".htm",
}


def _load_pdf(file_path: str) -> list[Document]:
    return PyPDFLoader(file_path).load()


def _load_csv(file_path: str) -> list[Document]:
    return CSVLoader(file_path).load()


def _load_docx(file_path: str) -> list[Document]:
    doc = DocxDocument(file_path)
    parts: list[str] = []

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)

    for table_index, table in enumerate(doc.tables, start=1):
        rows: list[str] = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            parts.append(f"Table {table_index}:\n" + "\n".join(rows))

    text = "\n\n".join(parts).strip()
    return [Document(page_content=text or "", metadata={"source": file_path})]


def _load_pptx(file_path: str) -> list[Document]:
    presentation = Presentation(file_path)
    slide_texts: list[str] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        items: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    items.append(text)
        if items:
            slide_texts.append(f"Slide {slide_index}:\n" + "\n".join(items))

    text = "\n\n".join(slide_texts).strip()
    return [Document(page_content=text or "", metadata={"source": file_path})]


def _load_spreadsheet(file_path: str) -> list[Document]:
    frame = pd.read_excel(file_path, sheet_name=None)
    documents: list[Document] = []

    for sheet_name, sheet_frame in frame.items():
        text = sheet_frame.fillna("").to_string(index=False)
        documents.append(Document(
            page_content=f"Sheet: {sheet_name}\n{text}".strip(),
            metadata={"source": file_path, "sheet": sheet_name},
        ))

    return documents


def _load_plain_text(file_path: str) -> list[Document]:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as handle:
        text = handle.read().strip()

    return [Document(page_content=text, metadata={"source": file_path})]


def _load_generic_text(file_path: str) -> list[Document]:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as handle:
        text = handle.read().strip()

    return [Document(page_content=text, metadata={"source": file_path})]


def load_and_chunk(file_path: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> list:
    """Load supported files and split them into chunks."""

    suffix = Path(file_path).suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file type: {file_path}")

    if suffix == ".pdf":
        documents = _load_pdf(file_path)
    elif suffix == ".csv":
        documents = _load_csv(file_path)
    elif suffix == ".docx":
        documents = _load_docx(file_path)
    elif suffix == ".pptx":
        documents = _load_pptx(file_path)
    elif suffix in {".xlsx", ".xls"}:
        documents = _load_spreadsheet(file_path)
    elif suffix in {".txt", ".md", ".rtf", ".json", ".html", ".htm"}:
        documents = _load_plain_text(file_path)
    else:
        documents = _load_generic_text(file_path)

    for doc in documents:
        doc.metadata["source"] = file_path

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap
    )

    chunks = splitter.split_documents(documents)
    return chunks